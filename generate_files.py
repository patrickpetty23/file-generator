#!/usr/bin/env python3
"""
Random File Generator
Generates various types of files with random content for testing purposes.
"""

import os
import random
import string
import json
import xml.etree.ElementTree as ET
import sqlite3
import wave
import struct
import zipfile
import csv
import io
import time
from pathlib import Path
from datetime import datetime, timedelta
from dotenv import load_dotenv

# Third-party imports (install via requirements.txt)
from PIL import Image, ImageDraw, ImageFont
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.colors import HexColor
from openpyxl import Workbook
from openpyxl.styles import Font, Fill, PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
import pandas as pd


# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================

def random_string(min_len=5, max_len=20):
    """Generate a random string of letters."""
    length = random.randint(min_len, max_len)
    return ''.join(random.choices(string.ascii_letters, k=length))


def random_word():
    """Generate a random word-like string."""
    return random_string(3, 12).lower()


def random_sentence(min_words=5, max_words=20):
    """Generate a random sentence."""
    num_words = random.randint(min_words, max_words)
    words = [random_word() for _ in range(num_words)]
    words[0] = words[0].capitalize()
    return ' '.join(words) + random.choice(['.', '!', '?'])


def random_paragraph(min_sentences=3, max_sentences=10):
    """Generate a random paragraph."""
    num_sentences = random.randint(min_sentences, max_sentences)
    return ' '.join(random_sentence() for _ in range(num_sentences))


def random_filename():
    """Generate a random filename (without extension)."""
    styles = [
        lambda: ''.join(random.choices(string.ascii_lowercase + string.digits, k=random.randint(6, 12))),
        lambda: f"{random_word()}_{random_word()}",
        lambda: f"{random_word()}-{random.randint(1, 9999)}",
        lambda: f"{random_word()}_{datetime.now().strftime('%Y%m%d')}_{random.randint(1, 999)}",
        lambda: ''.join(random.choices(string.hexdigits.lower(), k=8)),
    ]
    return random.choice(styles)()


def random_color():
    """Generate a random RGB color tuple."""
    return (random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))


def random_hex_color():
    """Generate a random hex color string."""
    return "#{:06x}".format(random.randint(0, 0xFFFFFF))


def random_date(start_year=2020, end_year=2026):
    """Generate a random date."""
    start = datetime(start_year, 1, 1)
    end = datetime(end_year, 12, 31)
    delta = end - start
    random_days = random.randint(0, delta.days)
    return start + timedelta(days=random_days)


def ensure_size_limit(content, max_bytes):
    """Ensure content doesn't exceed max bytes."""
    if isinstance(content, str):
        content = content.encode('utf-8')
    return content[:max_bytes]


# =============================================================================
# FILE GENERATORS
# =============================================================================

def generate_txt(filepath, max_size_mb=1):
    """Generate a random text file."""
    max_bytes = int(max_size_mb * 1024 * 1024)
    content = []
    current_size = 0
    
    while current_size < max_bytes:
        paragraph = random_paragraph() + "\n\n"
        content.append(paragraph)
        current_size += len(paragraph.encode('utf-8'))
        if current_size > max_bytes:
            break
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(''.join(content)[:max_bytes])


def generate_csv(filepath, max_size_mb=1):
    """Generate a random CSV file."""
    max_bytes = int(max_size_mb * 1024 * 1024)
    
    # Random column names
    num_cols = random.randint(3, 10)
    headers = [random_word() for _ in range(num_cols)]
    
    rows = []
    current_size = 0
    
    while current_size < max_bytes:
        row = []
        for _ in range(num_cols):
            cell_type = random.choice(['string', 'int', 'float', 'date'])
            if cell_type == 'string':
                row.append(random_word())
            elif cell_type == 'int':
                row.append(str(random.randint(-10000, 10000)))
            elif cell_type == 'float':
                row.append(f"{random.uniform(-1000, 1000):.2f}")
            else:
                row.append(random_date().strftime('%Y-%m-%d'))
        
        row_str = ','.join(row) + '\n'
        current_size += len(row_str.encode('utf-8'))
        rows.append(row)
        
        if current_size > max_bytes:
            break
    
    with open(filepath, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        writer.writerow(headers)
        writer.writerows(rows)


def generate_json(filepath, max_size_mb=1):
    """Generate a random JSON file."""
    max_bytes = int(max_size_mb * 1024 * 1024)
    
    def random_value(depth=0):
        if depth > 3:
            return random_word()
        
        value_type = random.choice(['string', 'int', 'float', 'bool', 'list', 'dict', 'null'])
        if value_type == 'string':
            return random_sentence()
        elif value_type == 'int':
            return random.randint(-10000, 10000)
        elif value_type == 'float':
            return round(random.uniform(-1000, 1000), 2)
        elif value_type == 'bool':
            return random.choice([True, False])
        elif value_type == 'null':
            return None
        elif value_type == 'list':
            return [random_value(depth + 1) for _ in range(random.randint(1, 5))]
        else:  # dict
            return {random_word(): random_value(depth + 1) for _ in range(random.randint(1, 5))}
    
    data = {random_word(): random_value() for _ in range(random.randint(5, 20))}
    
    content = json.dumps(data, indent=2)
    if len(content.encode('utf-8')) > max_bytes:
        content = content[:max_bytes]
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(content)


def generate_xml(filepath, max_size_mb=1):
    """Generate a random XML file."""
    max_bytes = int(max_size_mb * 1024 * 1024)
    
    root = ET.Element(random_word())
    
    def add_children(parent, depth=0, current_size=[0]):
        if depth > 4 or current_size[0] > max_bytes:
            return
        
        num_children = random.randint(1, 5)
        for _ in range(num_children):
            if current_size[0] > max_bytes:
                break
            child = ET.SubElement(parent, random_word())
            
            # Add attributes
            if random.random() > 0.5:
                child.set(random_word(), random_word())
            
            # Add text or more children
            if random.random() > 0.3 and depth < 4:
                add_children(child, depth + 1, current_size)
            else:
                child.text = random_sentence()
                current_size[0] += len(child.text.encode('utf-8'))
    
    add_children(root)
    
    tree = ET.ElementTree(root)
    tree.write(filepath, encoding='utf-8', xml_declaration=True)


def generate_html(filepath, max_size_mb=1):
    """Generate a random HTML file."""
    max_bytes = int(max_size_mb * 1024 * 1024)
    
    title = random_sentence(2, 5)
    
    html_parts = [
        '<!DOCTYPE html>',
        '<html lang="en">',
        '<head>',
        f'    <meta charset="UTF-8">',
        f'    <title>{title}</title>',
        f'    <style>',
        f'        body {{ font-family: {random.choice(["Arial", "Helvetica", "Georgia", "Times New Roman"])}; background-color: {random_hex_color()}; }}',
        f'        h1 {{ color: {random_hex_color()}; }}',
        f'        p {{ color: {random_hex_color()}; }}',
        f'    </style>',
        '</head>',
        '<body>',
    ]
    
    current_size = sum(len(p.encode('utf-8')) for p in html_parts)
    
    while current_size < max_bytes * 0.9:
        element_type = random.choice(['h1', 'h2', 'h3', 'p', 'div', 'ul', 'table'])
        
        if element_type in ['h1', 'h2', 'h3']:
            content = f'    <{element_type}>{random_sentence(3, 8)}</{element_type}>'
        elif element_type == 'p':
            content = f'    <p>{random_paragraph()}</p>'
        elif element_type == 'div':
            content = f'    <div style="background-color: {random_hex_color()}; padding: 10px;">{random_paragraph()}</div>'
        elif element_type == 'ul':
            items = ''.join(f'<li>{random_sentence()}</li>' for _ in range(random.randint(2, 6)))
            content = f'    <ul>{items}</ul>'
        else:  # table
            rows = random.randint(2, 5)
            cols = random.randint(2, 4)
            table_rows = []
            for i in range(rows):
                cells = ''.join(f'<td>{random_word()}</td>' for _ in range(cols))
                table_rows.append(f'<tr>{cells}</tr>')
            content = f'    <table border="1">{"".join(table_rows)}</table>'
        
        html_parts.append(content)
        current_size += len(content.encode('utf-8'))
        
        if current_size > max_bytes:
            break
    
    html_parts.extend(['</body>', '</html>'])
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write('\n'.join(html_parts))


def generate_md(filepath, max_size_mb=1):
    """Generate a random Markdown file."""
    max_bytes = int(max_size_mb * 1024 * 1024)
    
    content_parts = [f"# {random_sentence(3, 6)}\n\n"]
    current_size = len(content_parts[0].encode('utf-8'))
    
    while current_size < max_bytes:
        element_type = random.choice(['h2', 'h3', 'paragraph', 'list', 'code', 'quote', 'link'])
        
        if element_type == 'h2':
            content = f"## {random_sentence(2, 5)}\n\n"
        elif element_type == 'h3':
            content = f"### {random_sentence(2, 4)}\n\n"
        elif element_type == 'paragraph':
            content = f"{random_paragraph()}\n\n"
        elif element_type == 'list':
            items = '\n'.join(f"- {random_sentence()}" for _ in range(random.randint(2, 6)))
            content = f"{items}\n\n"
        elif element_type == 'code':
            code_content = '\n'.join(f"    {random_word()} = {random.randint(1, 100)}" for _ in range(random.randint(2, 5)))
            content = f"```\n{code_content}\n```\n\n"
        elif element_type == 'quote':
            content = f"> {random_sentence()}\n\n"
        else:  # link
            content = f"[{random_word()}](https://{random_word()}.com/{random_word()})\n\n"
        
        content_parts.append(content)
        current_size += len(content.encode('utf-8'))
        
        if current_size > max_bytes:
            break
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(''.join(content_parts))


def generate_png(filepath, max_size_mb=1):
    """Generate a random PNG image."""
    # Calculate dimensions based on max size (rough estimate)
    max_pixels = int(max_size_mb * 1024 * 1024 / 4)  # Rough estimate for RGBA
    side = min(int(max_pixels ** 0.5), 4096)  # Cap at 4096x4096
    
    width = random.randint(min(200, side), side)
    height = random.randint(min(200, side), side)
    
    img = Image.new('RGB', (width, height), random_color())
    draw = ImageDraw.Draw(img)
    
    # Draw random shapes
    num_shapes = random.randint(5, 50)
    for _ in range(num_shapes):
        shape_type = random.choice(['rectangle', 'ellipse', 'line', 'polygon'])
        color = random_color()
        
        if shape_type == 'rectangle':
            x1, y1 = random.randint(0, width), random.randint(0, height)
            x2, y2 = random.randint(0, width), random.randint(0, height)
            draw.rectangle([min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2)], fill=color)
        elif shape_type == 'ellipse':
            x1, y1 = random.randint(0, width), random.randint(0, height)
            x2, y2 = random.randint(0, width), random.randint(0, height)
            draw.ellipse([min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2)], fill=color)
        elif shape_type == 'line':
            points = [(random.randint(0, width), random.randint(0, height)) for _ in range(2)]
            draw.line(points, fill=color, width=random.randint(1, 10))
        else:  # polygon
            num_points = random.randint(3, 8)
            points = [(random.randint(0, width), random.randint(0, height)) for _ in range(num_points)]
            draw.polygon(points, fill=color)
    
    # Add some random text
    if random.random() > 0.5:
        text = random_word().upper()
        x = random.randint(0, max(1, width - 100))
        y = random.randint(0, max(1, height - 50))
        draw.text((x, y), text, fill=random_color())
    
    img.save(filepath, 'PNG')


def generate_jpg(filepath, max_size_mb=1):
    """Generate a random JPEG image."""
    max_pixels = int(max_size_mb * 1024 * 1024 / 3)
    side = min(int(max_pixels ** 0.5), 4096)
    
    width = random.randint(min(200, side), side)
    height = random.randint(min(200, side), side)
    
    img = Image.new('RGB', (width, height), random_color())
    draw = ImageDraw.Draw(img)
    
    # Create gradient or pattern
    pattern_type = random.choice(['gradient', 'noise', 'shapes'])
    
    if pattern_type == 'gradient':
        for x in range(width):
            r = int(255 * x / width)
            g = int(255 * (1 - x / width))
            b = random.randint(0, 255)
            for y in range(height):
                img.putpixel((x, y), (r, g, int(b * y / height)))
    elif pattern_type == 'noise':
        for x in range(0, width, 2):
            for y in range(0, height, 2):
                color = random_color()
                img.putpixel((x, y), color)
                if x + 1 < width:
                    img.putpixel((x + 1, y), color)
                if y + 1 < height:
                    img.putpixel((x, y + 1), color)
                if x + 1 < width and y + 1 < height:
                    img.putpixel((x + 1, y + 1), color)
    else:  # shapes
        for _ in range(random.randint(10, 100)):
            x1, y1 = random.randint(0, width), random.randint(0, height)
            x2, y2 = random.randint(0, width), random.randint(0, height)
            draw.ellipse([min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2)], fill=random_color())
    
    img.save(filepath, 'JPEG', quality=random.randint(60, 95))


def generate_gif(filepath, max_size_mb=1):
    """Generate a random animated GIF."""
    width = random.randint(100, 500)
    height = random.randint(100, 500)
    num_frames = random.randint(3, 15)
    
    frames = []
    for _ in range(num_frames):
        img = Image.new('RGB', (width, height), random_color())
        draw = ImageDraw.Draw(img)
        
        # Draw random shapes on each frame
        for _ in range(random.randint(3, 15)):
            x1, y1 = random.randint(0, width), random.randint(0, height)
            x2, y2 = random.randint(0, width), random.randint(0, height)
            shape = random.choice(['rectangle', 'ellipse'])
            if shape == 'rectangle':
                draw.rectangle([min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2)], fill=random_color())
            else:
                draw.ellipse([min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2)], fill=random_color())
        
        frames.append(img)
    
    frames[0].save(
        filepath,
        save_all=True,
        append_images=frames[1:],
        duration=random.randint(100, 500),
        loop=0
    )


def generate_bmp(filepath, max_size_mb=1):
    """Generate a random BMP image."""
    max_pixels = int(max_size_mb * 1024 * 1024 / 3)
    side = min(int(max_pixels ** 0.5), 2048)
    
    width = random.randint(min(100, side), side)
    height = random.randint(min(100, side), side)
    
    img = Image.new('RGB', (width, height), random_color())
    draw = ImageDraw.Draw(img)
    
    # Draw stripes or checkerboard
    pattern = random.choice(['stripes', 'checkerboard', 'random'])
    
    if pattern == 'stripes':
        stripe_width = random.randint(5, 50)
        for i in range(0, width, stripe_width * 2):
            draw.rectangle([i, 0, i + stripe_width, height], fill=random_color())
    elif pattern == 'checkerboard':
        cell_size = random.randint(10, 50)
        color1, color2 = random_color(), random_color()
        for x in range(0, width, cell_size):
            for y in range(0, height, cell_size):
                color = color1 if ((x // cell_size) + (y // cell_size)) % 2 == 0 else color2
                draw.rectangle([x, y, x + cell_size, y + cell_size], fill=color)
    else:
        for _ in range(random.randint(20, 100)):
            x1, y1 = random.randint(0, width), random.randint(0, height)
            x2, y2 = random.randint(0, width), random.randint(0, height)
            draw.rectangle([min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2)], fill=random_color())
    
    img.save(filepath, 'BMP')


def generate_svg(filepath, max_size_mb=1):
    """Generate a random SVG file."""
    width = random.randint(200, 1000)
    height = random.randint(200, 1000)
    
    svg_parts = [
        f'<?xml version="1.0" encoding="UTF-8"?>',
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">',
        f'  <rect width="100%" height="100%" fill="{random_hex_color()}"/>',
    ]
    
    num_shapes = random.randint(10, 100)
    for _ in range(num_shapes):
        shape_type = random.choice(['rect', 'circle', 'ellipse', 'line', 'polygon', 'text'])
        
        if shape_type == 'rect':
            x, y = random.randint(0, width), random.randint(0, height)
            w, h = random.randint(10, 200), random.randint(10, 200)
            svg_parts.append(f'  <rect x="{x}" y="{y}" width="{w}" height="{h}" fill="{random_hex_color()}" opacity="{random.uniform(0.3, 1):.2f}"/>')
        elif shape_type == 'circle':
            cx, cy = random.randint(0, width), random.randint(0, height)
            r = random.randint(10, 100)
            svg_parts.append(f'  <circle cx="{cx}" cy="{cy}" r="{r}" fill="{random_hex_color()}" opacity="{random.uniform(0.3, 1):.2f}"/>')
        elif shape_type == 'ellipse':
            cx, cy = random.randint(0, width), random.randint(0, height)
            rx, ry = random.randint(10, 100), random.randint(10, 100)
            svg_parts.append(f'  <ellipse cx="{cx}" cy="{cy}" rx="{rx}" ry="{ry}" fill="{random_hex_color()}" opacity="{random.uniform(0.3, 1):.2f}"/>')
        elif shape_type == 'line':
            x1, y1 = random.randint(0, width), random.randint(0, height)
            x2, y2 = random.randint(0, width), random.randint(0, height)
            svg_parts.append(f'  <line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" stroke="{random_hex_color()}" stroke-width="{random.randint(1, 10)}"/>')
        elif shape_type == 'polygon':
            points = ' '.join(f"{random.randint(0, width)},{random.randint(0, height)}" for _ in range(random.randint(3, 8)))
            svg_parts.append(f'  <polygon points="{points}" fill="{random_hex_color()}" opacity="{random.uniform(0.3, 1):.2f}"/>')
        else:  # text
            x, y = random.randint(0, width), random.randint(20, height)
            font_size = random.randint(12, 48)
            svg_parts.append(f'  <text x="{x}" y="{y}" font-size="{font_size}" fill="{random_hex_color()}">{random_word()}</text>')
    
    svg_parts.append('</svg>')
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write('\n'.join(svg_parts))


def generate_pdf(filepath, max_size_mb=1):
    """Generate a random PDF file."""
    c = canvas.Canvas(filepath, pagesize=random.choice([letter, A4]))
    page_width, page_height = letter
    
    num_pages = random.randint(1, max(1, int(max_size_mb * 5)))
    
    for page in range(num_pages):
        # Background color
        c.setFillColor(HexColor(random_hex_color()))
        c.rect(0, 0, page_width, page_height, fill=1)
        
        # Title
        c.setFillColor(HexColor(random_hex_color()))
        c.setFont("Helvetica-Bold", random.randint(18, 36))
        c.drawString(50, page_height - 50, random_sentence(3, 6))
        
        # Random text blocks
        y_position = page_height - 100
        c.setFont("Helvetica", random.randint(10, 14))
        
        while y_position > 100:
            c.setFillColor(HexColor(random_hex_color()))
            text = random_sentence()
            c.drawString(50, y_position, text[:80])  # Truncate long sentences
            y_position -= random.randint(15, 30)
        
        # Random shapes
        for _ in range(random.randint(3, 15)):
            c.setFillColor(HexColor(random_hex_color()))
            shape = random.choice(['rect', 'circle', 'line'])
            if shape == 'rect':
                x, y = random.randint(0, int(page_width)), random.randint(0, int(page_height))
                w, h = random.randint(20, 150), random.randint(20, 150)
                c.rect(x, y, w, h, fill=1)
            elif shape == 'circle':
                x, y = random.randint(0, int(page_width)), random.randint(0, int(page_height))
                r = random.randint(10, 75)
                c.circle(x, y, r, fill=1)
            else:  # line
                c.setStrokeColor(HexColor(random_hex_color()))
                c.setLineWidth(random.randint(1, 5))
                c.line(
                    random.randint(0, int(page_width)),
                    random.randint(0, int(page_height)),
                    random.randint(0, int(page_width)),
                    random.randint(0, int(page_height))
                )
        
        if page < num_pages - 1:
            c.showPage()
    
    c.save()


def generate_xlsx(filepath, max_size_mb=1):
    """Generate a random Excel file."""
    wb = Workbook()
    
    num_sheets = random.randint(1, 5)
    
    for sheet_idx in range(num_sheets):
        if sheet_idx == 0:
            ws = wb.active
            ws.title = random_word().capitalize()
        else:
            ws = wb.create_sheet(title=random_word().capitalize())
        
        num_rows = random.randint(10, 100)
        num_cols = random.randint(3, 15)
        
        # Header row
        for col in range(1, num_cols + 1):
            cell = ws.cell(row=1, column=col, value=random_word().capitalize())
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color=random_hex_color()[1:], end_color=random_hex_color()[1:], fill_type="solid")
        
        # Data rows
        for row in range(2, num_rows + 1):
            for col in range(1, num_cols + 1):
                cell_type = random.choice(['string', 'int', 'float', 'date', 'formula'])
                if cell_type == 'string':
                    ws.cell(row=row, column=col, value=random_word())
                elif cell_type == 'int':
                    ws.cell(row=row, column=col, value=random.randint(-10000, 10000))
                elif cell_type == 'float':
                    ws.cell(row=row, column=col, value=round(random.uniform(-1000, 1000), 2))
                elif cell_type == 'date':
                    ws.cell(row=row, column=col, value=random_date())
                else:  # formula
                    if row > 2:
                        ws.cell(row=row, column=col, value=f"=A{row-1}+1")
                    else:
                        ws.cell(row=row, column=col, value=random.randint(1, 100))
    
    wb.save(filepath)


def generate_pptx(filepath, max_size_mb=1):
    """Generate a random PowerPoint file."""
    prs = Presentation()
    
    num_slides = random.randint(3, 15)
    
    for _ in range(num_slides):
        slide_layout = prs.slide_layouts[random.randint(0, min(6, len(prs.slide_layouts) - 1))]
        slide = prs.slides.add_slide(slide_layout)
        
        # Try to set title if available
        if slide.shapes.title:
            slide.shapes.title.text = random_sentence(3, 7)
        
        # Add random shapes
        for _ in range(random.randint(2, 8)):
            shape_type = random.choice(['rectangle', 'oval', 'textbox'])
            
            left = Inches(random.uniform(0.5, 8))
            top = Inches(random.uniform(1, 6))
            width = Inches(random.uniform(0.5, 3))
            height = Inches(random.uniform(0.5, 2))
            
            if shape_type == 'rectangle':
                shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
            elif shape_type == 'oval':
                shape = slide.shapes.add_shape(9, left, top, width, height)  # MSO_SHAPE.OVAL
            else:  # textbox
                shape = slide.shapes.add_textbox(left, top, width, height)
                shape.text = random_sentence()
            
            if hasattr(shape, 'fill'):
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*random_color())
    
    prs.save(filepath)


def generate_docx(filepath, max_size_mb=1):
    """Generate a random Word document."""
    doc = Document()
    
    # Title
    doc.add_heading(random_sentence(3, 7), 0)
    
    num_sections = random.randint(5, 20)
    
    for _ in range(num_sections):
        section_type = random.choice(['heading', 'paragraph', 'list', 'table'])
        
        if section_type == 'heading':
            level = random.randint(1, 3)
            doc.add_heading(random_sentence(2, 5), level)
        elif section_type == 'paragraph':
            doc.add_paragraph(random_paragraph())
        elif section_type == 'list':
            for _ in range(random.randint(2, 6)):
                doc.add_paragraph(random_sentence(), style='List Bullet')
        else:  # table
            rows = random.randint(2, 6)
            cols = random.randint(2, 5)
            table = doc.add_table(rows=rows, cols=cols)
            table.style = 'Table Grid'
            for row in table.rows:
                for cell in row.cells:
                    cell.text = random_word()
    
    doc.save(filepath)


def generate_sqlite(filepath, max_size_mb=1):
    """Generate a random SQLite database."""
    if os.path.exists(filepath):
        os.remove(filepath)
    
    conn = sqlite3.connect(filepath)
    cursor = conn.cursor()
    
    num_tables = random.randint(2, 5)
    
    for table_idx in range(num_tables):
        table_name = f"{random_word()}_{table_idx}"
        num_cols = random.randint(3, 8)
        
        columns = []
        col_names = []
        for col_idx in range(num_cols):
            col_name = f"{random_word()}_{col_idx}"
            col_names.append(col_name)
            col_type = random.choice(['TEXT', 'INTEGER', 'REAL'])
            columns.append(f"{col_name} {col_type}")
        
        create_sql = f"CREATE TABLE {table_name} (id INTEGER PRIMARY KEY, {', '.join(columns)})"
        cursor.execute(create_sql)
        
        # Insert data
        num_rows = random.randint(50, 500)
        for _ in range(num_rows):
            values = []
            for col_name in col_names:
                val_type = random.choice(['string', 'int', 'float'])
                if val_type == 'string':
                    values.append(f"'{random_word()}'")
                elif val_type == 'int':
                    values.append(str(random.randint(-10000, 10000)))
                else:
                    values.append(str(round(random.uniform(-1000, 1000), 2)))
            
            insert_sql = f"INSERT INTO {table_name} ({', '.join(col_names)}) VALUES ({', '.join(values)})"
            cursor.execute(insert_sql)
    
    conn.commit()
    conn.close()


def generate_parquet(filepath, max_size_mb=1):
    """Generate a random Parquet file."""
    num_rows = random.randint(100, 10000)
    num_cols = random.randint(3, 10)
    
    data = {}
    for col_idx in range(num_cols):
        col_name = f"{random_word()}_{col_idx}"
        col_type = random.choice(['string', 'int', 'float', 'bool'])
        
        if col_type == 'string':
            data[col_name] = [random_word() for _ in range(num_rows)]
        elif col_type == 'int':
            data[col_name] = [random.randint(-10000, 10000) for _ in range(num_rows)]
        elif col_type == 'float':
            data[col_name] = [round(random.uniform(-1000, 1000), 2) for _ in range(num_rows)]
        else:
            data[col_name] = [random.choice([True, False]) for _ in range(num_rows)]
    
    df = pd.DataFrame(data)
    df.to_parquet(filepath, index=False)


def generate_wav(filepath, max_size_mb=1):
    """Generate a random WAV audio file."""
    # Parameters
    sample_rate = 44100
    max_samples = int(max_size_mb * 1024 * 1024 / 2)  # 2 bytes per sample
    duration = min(random.uniform(1, 30), max_samples / sample_rate)
    num_samples = int(sample_rate * duration)
    
    # Generate audio
    audio_type = random.choice(['sine', 'noise', 'mixed'])
    
    with wave.open(filepath, 'w') as wav_file:
        wav_file.setnchannels(1)  # Mono
        wav_file.setsampwidth(2)  # 2 bytes per sample
        wav_file.setframerate(sample_rate)
        
        for i in range(num_samples):
            if audio_type == 'sine':
                frequency = random.uniform(200, 2000)
                value = int(32767 * 0.5 * (1 + (i * frequency / sample_rate % 1)))
            elif audio_type == 'noise':
                value = random.randint(-32767, 32767)
            else:  # mixed
                frequency = random.uniform(200, 1000)
                sine_val = 0.5 * (1 + (i * frequency / sample_rate % 1))
                noise_val = random.uniform(-0.3, 0.3)
                value = int(32767 * (sine_val + noise_val))
                value = max(-32767, min(32767, value))
            
            wav_file.writeframes(struct.pack('<h', value))


def generate_zip(filepath, max_size_mb=1):
    """Generate a random ZIP archive."""
    with zipfile.ZipFile(filepath, 'w', zipfile.ZIP_DEFLATED) as zf:
        num_files = random.randint(3, 15)
        
        for _ in range(num_files):
            file_type = random.choice(['txt', 'json', 'csv'])
            inner_filename = f"{random_filename()}.{file_type}"
            
            if file_type == 'txt':
                content = random_paragraph()
            elif file_type == 'json':
                content = json.dumps({random_word(): random_sentence() for _ in range(random.randint(2, 10))})
            else:  # csv
                rows = [','.join(random_word() for _ in range(random.randint(3, 6))) for _ in range(random.randint(5, 20))]
                content = '\n'.join(rows)
            
            zf.writestr(inner_filename, content)


def generate_log(filepath, max_size_mb=1):
    """Generate a random log file."""
    max_bytes = int(max_size_mb * 1024 * 1024)
    
    log_levels = ['DEBUG', 'INFO', 'WARNING', 'ERROR', 'CRITICAL']
    modules = [random_word() for _ in range(5)]
    
    content = []
    current_size = 0
    
    current_time = datetime.now() - timedelta(days=random.randint(1, 30))
    
    while current_size < max_bytes:
        current_time += timedelta(seconds=random.randint(1, 300))
        timestamp = current_time.strftime('%Y-%m-%d %H:%M:%S')
        level = random.choice(log_levels)
        module = random.choice(modules)
        message = random_sentence()
        
        line = f"[{timestamp}] {level:8} {module}: {message}\n"
        content.append(line)
        current_size += len(line.encode('utf-8'))
        
        if current_size > max_bytes:
            break
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(''.join(content))


def generate_yaml(filepath, max_size_mb=1):
    """Generate a random YAML file."""
    def random_yaml_value(depth=0):
        if depth > 3:
            return random_word()
        
        value_type = random.choice(['string', 'int', 'float', 'bool', 'list', 'dict'])
        if value_type == 'string':
            return random_sentence()
        elif value_type == 'int':
            return random.randint(-1000, 1000)
        elif value_type == 'float':
            return round(random.uniform(-1000, 1000), 2)
        elif value_type == 'bool':
            return random.choice([True, False])
        elif value_type == 'list':
            return [random_yaml_value(depth + 1) for _ in range(random.randint(1, 4))]
        else:
            return {random_word(): random_yaml_value(depth + 1) for _ in range(random.randint(1, 4))}
    
    # Build YAML manually to avoid dependency
    def to_yaml(data, indent=0):
        lines = []
        prefix = "  " * indent
        
        if isinstance(data, dict):
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    lines.append(f"{prefix}{key}:")
                    lines.extend(to_yaml(value, indent + 1))
                else:
                    lines.append(f"{prefix}{key}: {value}")
        elif isinstance(data, list):
            for item in data:
                if isinstance(item, (dict, list)):
                    lines.append(f"{prefix}-")
                    lines.extend(to_yaml(item, indent + 1))
                else:
                    lines.append(f"{prefix}- {item}")
        else:
            lines.append(f"{prefix}{data}")
        
        return lines
    
    data = {random_word(): random_yaml_value() for _ in range(random.randint(5, 15))}
    content = '\n'.join(to_yaml(data))
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(content)


def generate_ini(filepath, max_size_mb=1):
    """Generate a random INI configuration file."""
    content = []
    
    num_sections = random.randint(3, 10)
    
    for _ in range(num_sections):
        section_name = random_word()
        content.append(f"[{section_name}]")
        
        num_options = random.randint(2, 8)
        for _ in range(num_options):
            key = random_word()
            value_type = random.choice(['string', 'int', 'float', 'bool', 'path'])
            
            if value_type == 'string':
                value = random_word()
            elif value_type == 'int':
                value = str(random.randint(0, 10000))
            elif value_type == 'float':
                value = f"{random.uniform(0, 100):.2f}"
            elif value_type == 'bool':
                value = random.choice(['true', 'false', 'yes', 'no', '1', '0'])
            else:  # path
                value = f"/path/to/{random_word()}/{random_word()}"
            
            content.append(f"{key} = {value}")
        
        content.append("")  # Empty line between sections
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write('\n'.join(content))


def generate_rtf(filepath, max_size_mb=1):
    """Generate a random RTF file."""
    content = [r"{\rtf1\ansi\deff0"]
    
    # Color table
    colors = [random_color() for _ in range(5)]
    color_table = r"{\colortbl;" + "".join(f"\\red{c[0]}\\green{c[1]}\\blue{c[2]};" for c in colors) + "}"
    content.append(color_table)
    
    num_paragraphs = random.randint(5, 30)
    
    for _ in range(num_paragraphs):
        color_idx = random.randint(1, len(colors))
        font_size = random.randint(20, 48)  # RTF font size is in half-points
        
        text = random_paragraph()
        formatted_text = f"\\cf{color_idx}\\fs{font_size} {text}\\par "
        content.append(formatted_text)
    
    content.append("}")
    
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(''.join(content))


# =============================================================================
# MAIN GENERATOR
# =============================================================================

# All available file generators
FILE_GENERATORS = {
    'txt': generate_txt,
    'csv': generate_csv,
    'json': generate_json,
    'xml': generate_xml,
    'html': generate_html,
    'md': generate_md,
    'png': generate_png,
    'jpg': generate_jpg,
    'gif': generate_gif,
    'bmp': generate_bmp,
    'svg': generate_svg,
    'pdf': generate_pdf,
    'xlsx': generate_xlsx,
    'pptx': generate_pptx,
    'docx': generate_docx,
    'sqlite': generate_sqlite,
    'parquet': generate_parquet,
    'wav': generate_wav,
    'zip': generate_zip,
    'log': generate_log,
    'yaml': generate_yaml,
    'ini': generate_ini,
    'rtf': generate_rtf,
}


def generate_random_file(output_dir, max_size_mb=1):
    """Generate a single random file of a random type."""
    file_type = random.choice(list(FILE_GENERATORS.keys()))
    filename = f"{random_filename()}.{file_type}"
    filepath = os.path.join(output_dir, filename)
    
    generator = FILE_GENERATORS[file_type]
    generator(filepath, max_size_mb)
    
    return filepath, file_type


def main():
    """Main function to generate random files based on .env configuration."""
    # Start timing
    start_time = time.time()
    
    # Load environment variables
    load_dotenv()
    
    # Get configuration
    output_path = os.getenv('OUTPUT_PATH', './generated_files')
    num_files = int(os.getenv('NUM_FILES', '10'))
    max_size_mb = float(os.getenv('MAX_FILE_SIZE_MB', '1'))
    
    # Validate max size
    if max_size_mb > 100:
        print(f"Warning: MAX_FILE_SIZE_MB ({max_size_mb}) exceeds 100MB cap. Using 100MB.")
        max_size_mb = 100
    
    # Create output directory
    output_dir = Path(output_path).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"Random File Generator")
    print(f"=" * 50)
    print(f"Output directory: {output_dir}")
    print(f"Number of files: {num_files}")
    print(f"Max file size: {max_size_mb} MB")
    print(f"Available file types: {len(FILE_GENERATORS)}")
    print(f"=" * 50)
    print()
    
    # Generate files
    generated = []
    type_counts = {}
    
    for i in range(num_files):
        try:
            filepath, file_type = generate_random_file(str(output_dir), max_size_mb)
            file_size = os.path.getsize(filepath)
            generated.append((filepath, file_type, file_size))
            type_counts[file_type] = type_counts.get(file_type, 0) + 1
            
            size_str = f"{file_size / 1024:.1f} KB" if file_size < 1024 * 1024 else f"{file_size / (1024 * 1024):.2f} MB"
            print(f"[{i + 1}/{num_files}] Generated: {os.path.basename(filepath)} ({file_type}, {size_str})")
        except Exception as e:
            print(f"[{i + 1}/{num_files}] Error generating file: {e}")
    
    # End timing
    end_time = time.time()
    elapsed_time = end_time - start_time
    
    # Summary
    print()
    print(f"=" * 50)
    print(f"Generation complete!")
    print(f"Total files generated: {len(generated)}")
    print(f"File type distribution:")
    for file_type, count in sorted(type_counts.items()):
        print(f"  {file_type}: {count}")
    
    total_size = sum(f[2] for f in generated)
    size_str = f"{total_size / 1024:.1f} KB" if total_size < 1024 * 1024 else f"{total_size / (1024 * 1024):.2f} MB"
    print(f"Total size: {size_str}")
    
    # Format elapsed time
    if elapsed_time < 1:
        time_str = f"{elapsed_time * 1000:.0f} ms"
    elif elapsed_time < 60:
        time_str = f"{elapsed_time:.2f} seconds"
    else:
        minutes = int(elapsed_time // 60)
        seconds = elapsed_time % 60
        time_str = f"{minutes}m {seconds:.2f}s"
    
    print(f"Time elapsed: {time_str}")


if __name__ == "__main__":
    main()
